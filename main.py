import os
import re
import shutil
from pptx import Presentation
import json

try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False


def load_rules():
    try:
        with open("rules.json", "r") as file:
            return json.load(file)
    except Exception as e:
        print(f"Error loading rules: {e}")
        return {}

SUBJECT_RULES = load_rules()

def preview_sort(folder_path, files):
    print("\nPlanned file organisation:")

    for file in files:
        category, score = get_category(folder_path, file)
        destination_folder = os.path.join(folder_path, category)
        print(f"Would move '{file}' -> '{destination_folder}' (score: {score})")

def create_folder_if_needed(folder_path):
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)

def get_unique_destination(destination_path):
    if not os.path.exists(destination_path):
        return destination_path

    base, ext = os.path.splitext(destination_path)
    counter = 1

    while True:
        new_path = f"{base}_{counter}{ext}"
        if not os.path.exists(new_path):
            return new_path
        counter += 1

def sort_files(folder_path, files):
    print("\nSorting files...")

    for file in files:
        source_path = os.path.join(folder_path, file)

        if not os.path.exists(source_path):
            continue

        category, score = get_category(folder_path, file)

        destination_folder = os.path.join(folder_path, category)
        create_folder_if_needed(destination_folder)

        destination_path = os.path.join(destination_folder, file)
        destination_path = get_unique_destination(destination_path)

        try:
            shutil.move(source_path, destination_path)
            print(f"Moved '{file}' -> '{destination_folder}' (score: {score})")
        except Exception as e:
            print(f"Could not move '{file}': {e}")


def get_files(folder_path):
    try:
        items = os.listdir(folder_path)
        files = []

        for item in items:
            # skip hidden files (mac junk)
            if item.startswith("."):
                continue

            full_path = os.path.join(folder_path, item)

            # skip if not a real file
            if not os.path.isfile(full_path):
                continue

            # skip mac apps and weird system stuff
            if item.endswith(".app"):
                continue

            files.append(item)

        return files

    except FileNotFoundError:
        print("That folder does not exist.")
        return []
    except Exception as e:
        print(f"An error occurred: {e}")
        return []


def clean_text(text):
    text = text.lower()
    text = re.sub(r"[^a-z0-9]+", " ", text)
    return text


def tokenize(text):
    return re.findall(r"[a-z]+[0-9]*|[0-9]+", text.lower())


def read_text_file(full_path):
    try:
        with open(full_path, "r", encoding="utf-8", errors="ignore") as file:
            return file.read()
    except Exception:
        return ""


def read_pdf_file(full_path):
    if not PDF_AVAILABLE:
        return ""

    try:
        reader = PdfReader(full_path)
        text = ""

        for page in reader.pages:
            try:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            except Exception:
                continue

        return text
    except Exception:
        return ""

def read_pptx_file(full_path):
    try:
        prs = Presentation(full_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text
    except Exception:
        return ""

def extract_file_text(full_path):
    ext = os.path.splitext(full_path)[1].lower()

    if ext in [".txt", ".py", ".md", ".csv", ".json", ".html", ".css", ".js", ".ts", ".sql"]:
        return read_text_file(full_path)

    if ext == ".pdf":
        return read_pdf_file(full_path)

    if ext == ".pptx":
        return read_pptx_file(full_path)

    return ""

def score_subject_from_text(text, rules):
    cleaned = clean_text(text)
    tokens = tokenize(cleaned)
    score = 0

    for keyword in rules["strong"]:
        if keyword in cleaned:
            score += 10

    for keyword in rules["medium"]:
        if keyword in tokens:
            score += 4
        elif keyword in cleaned:
            score += 2

    for keyword in rules["weak"]:
        if keyword in tokens:
            score += 2
        elif keyword in cleaned:
            score += 1

    return score


def get_subject_from_content(file_text):
    best_subject = "Needs_Review"
    best_score = 0

    for subject, rules in SUBJECT_RULES.items():
        score = score_subject_from_text(file_text, rules)

        if score > best_score:
            best_score = score
            best_subject = subject

    if best_score >= 8:
        return best_subject, best_score

    return "Needs_Review", best_score


def get_subject_from_filename(filename):
    cleaned = clean_text(os.path.splitext(filename)[0])
    best_subject = "Needs_Review"
    best_score = 0

    for subject, rules in SUBJECT_RULES.items():
        score = score_subject_from_text(cleaned, rules)

        if score > best_score:
            best_score = score
            best_subject = subject

    if best_score >= 5:
        return best_subject, best_score

    return "Needs_Review", best_score

def get_file_type_category(filename):
    ext = os.path.splitext(filename)[1].lower()

    if ext in [".jpg", ".jpeg", ".png", ".gif", ".heic"]:
        return "Images"
    elif ext in [".pdf",".ppt",".doc", ".docx", ".txt", ".pptx", ".xlsx", ".csv"]:
        return "Documents"
    elif ext in [".mp4", ".mov", ".avi"]:
        return "Videos"
    elif ext in [".zip", ".rar"]:
        return "Archives"
    elif ext in [".py", ".js", ".ts", ".java", ".html", ".css", ".sql"]:
        return "Code"
    else:
        return "Other"


def get_category(folder_path, filename):
    full_path = os.path.join(folder_path, filename)

    file_text = extract_file_text(full_path)
    if file_text.strip():
        subject, score = get_subject_from_content(file_text)
        if subject != "Needs_Review":
            return subject, score

    subject, score = get_subject_from_filename(filename)
    if subject != "Needs_Review":
        return subject, score

    fallback = get_file_type_category(filename)
    return fallback, score


folder = input("Enter folder path: ")
files = get_files(folder)

print("\nFiles found:")
for file in files:
    category, score = get_category(folder, file)
    print(f"{file} -> {category} (score: {score})")

preview_sort(folder, files)

confirm = input("\nDo you want to sort these files for real? (yes/no): ").strip().lower()

if confirm == "yes":
    sort_files(folder, files)
    print("\nDone.")
else:
    print("\nNo files were moved.")